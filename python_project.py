import sys
import pyttsx3
import requests
import threading
import pythoncom  # Required for COM initialization on Windows
from PyQt6.QtWidgets import QApplication, QWidget, QVBoxLayout, QPushButton, QTextEdit, QFileDialog, QLabel
from PyQt6.QtGui import QFont
from PyQt6.QtCore import QThread
from pptx import Presentation

# API Key for text summarization
API_KEY = "Enter your API Key"
API_URL = "Enter your API URL" 

class SpeechThread(QThread):
    def __init__(self, text, parent=None):
        super().__init__(parent)
        self.text = text
``
    def run(self):
        pythoncom.CoInitialize()
        try:
            engine = pyttsx3.init()
            engine.say(self.text)
            engine.runAndWait()
        finally:
            pythoncom.CoUninitialize()

class AITutorApp(QWidget):
    def __init__(self):
        super().__init__()
        self.initUI()
        self.full_text = ""  # Store extracted PPT text
        self.speech_thread = None

    def initUI(self):
        self.setWindowTitle("SummarEase - AI Tutor 🧠")
        self.setGeometry(100, 100, 600, 400)

        layout = QVBoxLayout()

        self.label = QLabel("Upload a PPT file:")
        self.label.setFont(QFont("Comic Sans MS", 15))
        layout.addWidget(self.label)

        self.uploadButton = QPushButton("Upload PPT")
        self.uploadButton.setFont(QFont("Comic Sans MS", 10))
        self.uploadButton.clicked.connect(self.upload_ppt)
        layout.addWidget(self.uploadButton)

        self.textArea = QTextEdit()
        self.textArea.setFont(QFont("Comic Sans MS", 10))
        self.textArea.setPlaceholderText("Lesson summary will appear here...")
        layout.addWidget(self.textArea)

        self.generateButton = QPushButton("Generate Summary")
        self.generateButton.setFont(QFont("Comic Sans MS", 10))
        self.generateButton.clicked.connect(self.generate_summary)
        layout.addWidget(self.generateButton)

        self.speakButton = QPushButton("Speak Summary")
        self.speakButton.setFont(QFont("Comic Sans MS", 10))
        self.speakButton.clicked.connect(self.speak_text)
        layout.addWidget(self.speakButton)

        self.setLayout(layout)

    def upload_ppt(self):
        file_name, _ = QFileDialog.getOpenFileName(self, "Open PPT File", "", "PowerPoint Files (*.pptx)")
        if file_name:
            self.extract_text_from_ppt(file_name)
            self.textArea.setText("PPT successfully uploaded! Now generate the summary.")
            self.run_speech("PowerPoint successfully uploaded. Now generate the summary.")

    def extract_text_from_ppt(self, file_name):
        prs = Presentation(file_name)
        self.full_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    self.full_text += shape.text + "\n"

    def generate_summary(self):
        if self.full_text.strip():
            summary = self.summarize_text(self.full_text)
            self.textArea.setText(summary)  # Display only the summary
            self.run_speech(summary)
        else:
            self.textArea.setText("No content to summarize!")
            self.run_speech("No content to summarize!")

    def summarize_text(self, text):
        headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
        data = {"inputs": text, "parameters": {"max_length": 300, "min_length": 100}}  # Adjust length as needed
        
        response = requests.post(API_URL, json=data, headers=headers)
        if response.status_code == 200:
            try:
                result = response.json()
                summary = result[0].get("summary_text", "Error: No summary generated")
                return summary
            except Exception as e:
                print("Error parsing JSON:", e)
                return "Error: Unable to parse summary response"
        else:
            # Debug output for troubleshooting
            print("Response Code:", response.status_code)
            print("Response Text:", response.text)
            return "Error: Unable to summarize text"

    def speak_text(self):
        summary = self.textArea.toPlainText().strip()
        if summary:
            self.run_speech(summary)
        else:
            self.textArea.append("\nNo summary to speak!")
            self.run_speech("No summary to speak!")

    def run_speech(self, text):
        # If there is an active speech thread, wait for it to finish
        if self.speech_thread and self.speech_thread.isRunning():
            self.speech_thread.wait()
        self.speech_thread = SpeechThread(text)
        self.speech_thread.start()

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = AITutorApp()
    window.show()
    sys.exit(app.exec())
